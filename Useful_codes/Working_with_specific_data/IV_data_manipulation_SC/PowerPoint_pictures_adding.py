import os
import sys
import time
from PIL import Image
from tqdm import tqdm
from pptx import Presentation
from pptx.util import Inches
from icecream import ic


class ImagePresentation:
    """
    A class to create a PowerPoint presentation with images from a given directory.

    :param images_directory: The directory containing the images to be added to the presentation.
    :type images_directory: str
    """
    def __init__(self, images_directory):
        """
        Initializes the ImagePresentation class with the given images_directory.

        :param images_directory: The directory containing the images to be added to the presentation.
        :type images_directory: str
        """
        self.images_directory = images_directory
        self.presentation = Presentation()
        self.image_extensions = ['.jpg', '.jpeg', '.png', '.gif', '.bmp', 'svg']

    def add_image_slide(self, image_path):
        """
        Adds an image to a new slide in the presentation.

        :param image_path: The path of the image to be added.
        :type image_path: str
        :return: The created slide with the added image.
        :rtype: Slide
        """
        slide_layout = self.presentation.slide_layouts[5]
        slide = self.presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = os.path.basename(image_path)
        img = Image.open(image_path)
        dpi = img.info.get('dpi', (300, 300))
        dpi_x, dpi_y = dpi

        width, height = img.size
        width_cm = width * (2.54 / dpi_x)
        height_cm = height * (2.54 / dpi_y)
        max_width_cm = 25
        max_height_cm = 15

        aspect_ratio = float(width_cm) / float(height_cm)
        if width_cm > max_width_cm:
            width_cm = max_width_cm
            height_cm = width_cm / aspect_ratio
        if height_cm > max_height_cm:
            height_cm = max_height_cm
            width_cm = height_cm * aspect_ratio

        file_size = os.stat(image_path).st_size
        if file_size > 2 * 1024 * 1024:  # Check if the file size is greater than 2 MB
            new_width = int(width_cm * (dpi_x / 2.54))
            new_height = int(height_cm * (dpi_y / 2.54))
            img = img.resize((new_width, new_height), Image.LANCZOS)
            img.save('temp_resized_image.png', 'PNG')
            image_path = 'temp_resized_image.png'

        slide_width = self.presentation.slide_width.inches
        slide_height = self.presentation.slide_height.inches

        title_height = 1
        slide_height -= title_height

        left = (slide_width - (width_cm * 0.393701)) / 2
        top = slide_height - (height_cm * 0.393701) + title_height
        slide.shapes.add_picture(image_path, Inches(left), Inches(top),
                                 width=Inches(width_cm * 0.393701), height=Inches(height_cm * 0.393701))

        if 'temp_resized_image.png' in image_path:
            os.remove('temp_resized_image.png')

        return slide

    def create_presentation(self):
        """
        Creates a PowerPoint presentation with the images found in the specified directory.
        """
        if not os.path.isdir(self.images_directory):
            print(f"Error: {self.images_directory} is not a directory.")
            sys.exit(1)

        image_files = [f for f in os.listdir(self.images_directory) if os.path.splitext(f)[1].lower() in self.image_extensions]

        if not image_files:
            print(f"Error: No image files found in {self.images_directory}.")
            sys.exit(1)

        for image_file in tqdm(image_files, desc="Adding images to slides"):
            image_path = os.path.join(self.images_directory, image_file)
            self.add_image_slide(image_path)

        output_presentation = os.path.join(self.images_directory, "output_presentation.pptx")
        self.presentation.save(output_presentation)
        print(f"Presentation saved as {output_presentation}")
        os.startfile(output_presentation)


if __name__ == '__main__':
    ic.configureOutput('')
    start_time = time.time()
    images_directory_path = r"D:\OneDrive - O365 Turun yliopisto\Desktop\Aralab\Revival_2_killing/delete me"
    image_presentation = ImagePresentation(images_directory_path)
    image_presentation.create_presentation()
    print("\n", "--- %s seconds ---" % (time.time() - start_time))